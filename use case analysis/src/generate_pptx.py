#!/usr/bin/env python3
"""
generate_pptx.py
────────────────
Génère la présentation PowerPoint (.pptx) Air Liquide — AI Builders
avec la charte graphique Pyl.Tech via python-pptx pur (sans PyltechDeck).

Architecture : fonctions build_* indépendantes + main() orchestrateur.
"""
import sys
from pathlib import Path
from datetime import date

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PROJECT_DIR  = Path(__file__).parent.parent
CATALOG_FILE = PROJECT_DIR / "output" / "use_cases_catalog.xlsx"
OUTPUT_FILE  = PROJECT_DIR / "output" / "presentation_ai_champions.pptx"

# ── Slide dimensions ──────────────────────────────────────────────────────────
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── Brand colors ──────────────────────────────────────────────────────────────
PYL_NAVY_DARK = RGBColor(0x0B, 0x13, 0x2B)
PYL_NAVY      = RGBColor(0x0D, 0x21, 0x49)
PYL_YELLOW    = RGBColor(0xF4, 0xBF, 0x46)
PYL_TEAL_BLUE = RGBColor(0x20, 0x8A, 0xAE)
PYL_TEAL      = RGBColor(0x5B, 0xC0, 0xBE)
PYL_BODY_GREY = RGBColor(0x4F, 0x4F, 0x4F)
PYL_GREY_BG   = RGBColor(0xEE, 0xEE, 0xEE)
PYL_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
PYL_SUCCESS   = RGBColor(0x13, 0x86, 0x36)
PYL_DANGER    = RGBColor(0xC9, 0x14, 0x32)


# ── Internal ──────────────────────────────────────────────────────────────────

def _blank_slide(prs: Presentation):
    """Retourne un slide vierge (layout index 6)."""
    return prs.slides.add_slide(prs.slide_layouts[6])


# ── Low-level drawing helpers ─────────────────────────────────────────────────

def add_rect(slide, left, top, width, height,
             fill: RGBColor = None, no_border: bool = True):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if no_border:
        shape.line.fill.background()
    return shape


def add_textbox(
    slide, left, top, width, height, text: str,
    font_name: str = "Poppins", font_size: int = 12,
    bold: bool = False, italic: bool = False,
    color: RGBColor = None, align=PP_ALIGN.LEFT, word_wrap: bool = True,
):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = word_wrap
    p = txBox.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return txBox


def add_line(slide, x1, y1, x2, y2,
             color: RGBColor = None, width_pt: float = 1.0):
    """Ajoute un connecteur ligne entre (x1, y1) et (x2, y2)."""
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)
    connector.line.width = Pt(width_pt)
    if color:
        connector.line.color.rgb = color
    return connector


def add_footer(slide, page: int = 0,
               label: str = "Air Liquide — AI Builders") -> None:
    """Barre de pied de page Pyl.Tech. Si page > 0, affiche le numéro."""
    add_rect(slide, Inches(0), Inches(7.1), SLIDE_W, Inches(0.4), fill=PYL_NAVY_DARK)
    add_textbox(
        slide, Inches(0.3), Inches(7.15), Inches(10), Inches(0.3),
        label, font_size=8, color=PYL_WHITE,
    )
    if page > 0:
        add_textbox(
            slide, Inches(12.5), Inches(7.15), Inches(0.8), Inches(0.3),
            str(page), font_size=8, bold=True, color=PYL_YELLOW,
            align=PP_ALIGN.RIGHT,
        )


def add_slide_header(slide, chapter_num: str, title: str,
                     subtitle: str = "", page: int = 0) -> None:
    """Bande navy + numéro de chapitre + titre principal."""
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), fill=PYL_NAVY)
    add_textbox(
        slide, Inches(0.3), Inches(0.1), Inches(1.0), Inches(0.4),
        chapter_num, font_size=28, bold=True, color=PYL_YELLOW,
    )
    add_textbox(
        slide, Inches(1.5), Inches(0.15), Inches(10.5), Inches(0.5),
        title, font_size=20, bold=True, color=PYL_WHITE,
    )
    if subtitle:
        add_textbox(
            slide, Inches(0.3), Inches(0.75), Inches(12.0), Inches(0.35),
            subtitle, font_size=11, italic=True, color=PYL_TEAL_BLUE,
        )
    add_footer(slide, page=page)


def add_kpi_card(
    slide, left, top, width, height,
    label: str, value: str, sub: str = "",
    header_color: RGBColor = PYL_NAVY,
    value_color: RGBColor = PYL_TEAL_BLUE,
) -> None:
    card = add_rect(slide, left, top, width, height, fill=PYL_WHITE)
    card.line.color.rgb = PYL_GREY_BG
    card.line.width = Pt(0.5)

    add_rect(slide, left, top, width, Inches(0.07), fill=header_color)

    v_box = slide.shapes.add_textbox(
        left + Inches(0.1), top + Inches(0.15), width - Inches(0.2), Inches(0.9))
    v_box.text_frame.word_wrap = False
    p = v_box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = value
    run.font.name = "Poppins"
    run.font.bold = True
    run.font.size = Pt(36)
    run.font.color.rgb = value_color

    l_box = slide.shapes.add_textbox(
        left + Inches(0.1), top + Inches(1.0), width - Inches(0.2), Inches(0.5))
    l_box.text_frame.word_wrap = True
    p = l_box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.name = "Poppins"
    run.font.size = Pt(10)
    run.font.color.rgb = PYL_BODY_GREY

    if sub:
        s_box = slide.shapes.add_textbox(
            left + Inches(0.1), top + Inches(1.45), width - Inches(0.2), Inches(0.35))
        s_box.text_frame.word_wrap = True
        p = s_box.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = sub
        run.font.name = "Poppins"
        run.font.size = Pt(9)
        run.font.italic = True
        run.font.color.rgb = PYL_BODY_GREY


def add_content_card(
    slide, left, top, width, height,
    header_color: RGBColor, card_title: str, bullets: list,
    body_font_size: int = 11,
) -> None:
    card_bg = add_rect(slide, left, top, width, height, fill=PYL_WHITE)
    card_bg.line.color.rgb = PYL_GREY_BG
    card_bg.line.width = Pt(0.5)

    add_rect(slide, left, top, width, Inches(0.45), fill=header_color)

    t_box = slide.shapes.add_textbox(
        left + Inches(0.15), top + Inches(0.06), width - Inches(0.3), Inches(0.35))
    p = t_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = card_title
    run.font.name = "Poppins"
    run.font.bold = True
    run.font.size = Pt(13)
    run.font.color.rgb = PYL_WHITE

    line_h = Inches(0.42) if body_font_size <= 11 else Inches(0.78)
    cur_top = top + Inches(0.55)
    for bullet in bullets:
        b_box = slide.shapes.add_textbox(
            left + Inches(0.15), cur_top, width - Inches(0.3), Inches(0.45))
        b_box.text_frame.word_wrap = True
        p = b_box.text_frame.paragraphs[0]
        sq = p.add_run()
        sq.text = "  "
        sq.font.name = "Poppins"
        sq.font.bold = True
        sq.font.size = Pt(body_font_size)
        sq.font.color.rgb = PYL_YELLOW
        txt = p.add_run()
        txt.text = bullet
        txt.font.name = "Poppins"
        txt.font.size = Pt(body_font_size)
        txt.font.color.rgb = PYL_BODY_GREY
        cur_top += line_h


def add_pptx_table(
    slide, left, top, width, height,
    headers: list, rows: list,
) -> None:
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table

    col_w = width // n_cols
    for i in range(n_cols):
        tbl.columns[i].width = col_w

    for i, h in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PYL_NAVY_DARK
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = h
        run.font.name = "Poppins"
        run.font.bold = True
        run.font.size = Pt(10)
        run.font.color.rgb = PYL_WHITE

    for r_idx, row_data in enumerate(rows):
        bg = PYL_NAVY if r_idx % 2 == 0 else PYL_WHITE
        fg = PYL_WHITE if r_idx % 2 == 0 else PYL_BODY_GREY
        for c_idx, val in enumerate(row_data):
            cell = tbl.cell(r_idx + 1, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if c_idx > 0 else PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(val)
            run.font.name = "Poppins"
            run.font.size = Pt(10)
            run.font.color.rgb = fg


# ── Slide builders ────────────────────────────────────────────────────────────

def build_cover(prs: Presentation) -> None:
    """Slide 01 — Couverture."""
    slide = _blank_slide(prs)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=PYL_NAVY)
    add_textbox(
        slide, Inches(0.5), Inches(2.5), Inches(12), Inches(1.2),
        "Air Liquide — AI Builders",
        font_size=36, bold=True, color=PYL_WHITE, align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide, Inches(0.5), Inches(3.8), Inches(12), Inches(0.6),
        "Analyse, Classification et Recommandations du Portefeuille Use Cases IA",
        font_size=16, italic=True, color=PYL_YELLOW, align=PP_ALIGN.CENTER,
    )
    add_rect(slide, Inches(3), Inches(4.8), Inches(7.33), Inches(0.05), fill=PYL_YELLOW)
    add_textbox(
        slide, Inches(0.5), Inches(5.0), Inches(12), Inches(0.4),
        date.today().strftime("%B %Y"),
        font_size=12, color=PYL_TEAL_BLUE, align=PP_ALIGN.CENTER,
    )


def build_kpi_slide(prs: Presentation, df: pd.DataFrame, page: int = 0) -> None:
    """Slide KPI — 6 cartes chiffres clés."""
    slide  = _blank_slide(prs)
    total  = len(df)
    tiers  = df["Complexity_Tier"].value_counts().to_dict()
    it_count = int((df["IT_Flag"] != "").sum())
    clusters = df["Cluster"].nunique()

    add_slide_header(slide, "02", f"{total} use cases analysés", page=page)

    CARD_W = Inches(1.5)
    CARD_H = Inches(1.85)
    Y_TOP  = Inches(1.75)
    GAP    = Inches(0.1)
    LEFT   = Inches(0.2)

    cards = [
        ("Use Cases",  str(total),                      "Portefeuille",  PYL_NAVY,    PYL_YELLOW),
        ("Clusters",   str(clusters),                   "Périmètre géo", PYL_NAVY,    PYL_TEAL_BLUE),
        ("Small",      str(tiers.get("Small", 0)),      "< 2 semaines",  PYL_SUCCESS, PYL_SUCCESS),
        ("Medium",     str(tiers.get("Medium", 0)),     "4-8 semaines",  PYL_TEAL,    PYL_TEAL),
        ("Large",      str(tiers.get("Large", 0)),      "3-12 mois",     PYL_DANGER,  PYL_DANGER),
        ("Points IT",  str(it_count),                   "Escalade",      PYL_DANGER,  PYL_DANGER),
    ]
    for i, (label, value, sub, hdr_color, val_color) in enumerate(cards):
        left = LEFT + i * (CARD_W + GAP)
        add_kpi_card(slide, left, Y_TOP, CARD_W, CARD_H, label, value, sub,
                     header_color=hdr_color, value_color=val_color)

    pct_accessible = (tiers.get("Small", 0) + tiers.get("Medium", 0)) * 100 // max(total, 1)
    add_textbox(
        slide, LEFT, Inches(3.75), Inches(9.6), Inches(0.38),
        f"{pct_accessible}% des use cases sont déployables sans ressources IT lourdes",
        font_size=11, bold=True, color=PYL_NAVY_DARK, align=PP_ALIGN.CENTER,
    )


def build_complexity_slide(prs: Presentation, df: pd.DataFrame, page: int = 0) -> None:
    """Slide Complexité — 3 cartes Small / Medium / Large."""
    slide = _blank_slide(prs)
    total = len(df)
    tiers = df["Complexity_Tier"].value_counts().to_dict()

    add_slide_header(slide, "02", "Scoring : Small / Medium / Large", page=page)

    tier_data = [
        ("Small",  tiers.get("Small", 0),  "Quick Win",           "< 2 sem.",  "1 builder, no-code",    PYL_SUCCESS),
        ("Medium", tiers.get("Medium", 0), "Use Case Structurant", "4-8 sem.",  "builder + IT local",    PYL_TEAL_BLUE),
        ("Large",  tiers.get("Large", 0),  "Projet Stratégique",  "3-12 mois", "Équipe IT + builder",   PYL_DANGER),
    ]
    CARD_W = Inches(2.95)
    CARD_H = Inches(3.55)
    Y_TOP  = Inches(1.75)

    for i, (tier, count, label, ttv, profile, color) in enumerate(tier_data):
        left = Inches(0.3) + i * (CARD_W + Inches(0.15))
        pct  = count * 100 // max(total, 1)
        add_content_card(
            slide, left, Y_TOP, CARD_W, CARD_H, color,
            f"{tier} — {count} UC ({pct}%)",
            [f"Label : {label}", f"Time-to-value : {ttv}", f"Profil : {profile}"],
        )
        add_textbox(
            slide, left, Y_TOP + Inches(0.55), CARD_W, Inches(0.8),
            str(count), font_size=48, bold=True, color=PYL_WHITE, align=PP_ALIGN.CENTER,
        )


def build_families_slide(prs: Presentation, df: pd.DataFrame, page: int = 0) -> None:
    """Slide Familles — tableau par famille fonctionnelle."""
    slide = _blank_slide(prs)
    add_slide_header(slide, "02", "7 familles fonctionnelles identifiées", page=page)

    families = df[["Family_Label", "Family"]].drop_duplicates("Family_Label").copy()
    families["_sort"] = families["Family"].str.extract(r"(\d+)").astype(int)
    families = families.sort_values("_sort")

    rows = []
    for _, row in families.iterrows():
        fam    = row["Family_Label"]
        code   = row["Family"]
        subset = df[df["Family_Label"] == fam]["Complexity_Tier"].value_counts().to_dict()
        rows.append([
            f"{code} — {fam}",
            str(subset.get("Small", 0)),
            str(subset.get("Medium", 0)),
            str(subset.get("Large", 0)),
        ])

    add_pptx_table(
        slide, Inches(0.3), Inches(1.75), Inches(9.4), Inches(3.6),
        ["Famille fonctionnelle", "Small", "Medium", "Large"], rows,
    )


def build_clusters_slide(prs: Presentation, df: pd.DataFrame, page: int = 0) -> None:
    """Slide Clusters — top 12 clusters par complexité."""
    slide = _blank_slide(prs)
    add_slide_header(slide, "02", "Top clusters", page=page)

    cluster_agg = (
        df.groupby("Cluster")["Complexity_Tier"]
        .value_counts()
        .unstack(fill_value=0)
        .reset_index()
    )
    # Garantit les 3 colonnes même si un tier est absent du jeu de données
    for col in ("Small", "Medium", "Large"):
        if col not in cluster_agg.columns:
            cluster_agg[col] = 0

    cluster_agg["Total"] = (
        cluster_agg["Small"] + cluster_agg["Medium"] + cluster_agg["Large"]
    )
    cluster_agg = cluster_agg.sort_values("Total", ascending=False).head(12)

    rows = []
    for _, row in cluster_agg.iterrows():
        rows.append([
            str(row["Cluster"]),
            str(int(row["Small"])),
            str(int(row["Medium"])),
            str(int(row["Large"])),
            str(int(row["Total"])),
        ])

    add_pptx_table(
        slide, Inches(0.3), Inches(1.75), Inches(9.4), Inches(3.6),
        ["Cluster", "Small", "Medium", "Large", "Total"], rows,
    )


def build_it_slide(prs: Presentation, df: pd.DataFrame, page: int = 0) -> None:
    """Slide Points IT — top UC avec dépendances systèmes."""
    slide = _blank_slide(prs)
    add_slide_header(slide, "03", "Points d'attention IT", page=page)

    it_df = df[df["IT_Flag"] != ""]
    if it_df.empty:
        add_textbox(
            slide, Inches(0.3), Inches(2.5), Inches(9), Inches(0.5),
            "Aucun point d'attention IT identifié.", font_size=12, color=PYL_BODY_GREY,
        )
        return

    top_it = it_df.nlargest(8, "Score_Total")
    rows = []
    for _, row in top_it.iterrows():
        raw  = str(row["Use Case Description (Long)"])
        desc = raw[:50] + ("..." if len(raw) > 50 else "")
        rows.append([str(row["UC_ID"]), str(row["Cluster"]), str(row["Complexity_Tier"]), desc])

    add_pptx_table(
        slide, Inches(0.3), Inches(1.75), Inches(9.4), Inches(3.6),
        ["UC_ID", "Cluster", "Tier", "Description"], rows,
    )


def build_quickwins_slide(prs: Presentation, df: pd.DataFrame, page: int = 0) -> None:
    """Slide Quick Wins — UC Small déployables immédiatement."""
    slide = _blank_slide(prs)
    add_slide_header(slide, "04", "Top Quick Wins", page=page)

    qw = df[(df["Complexity_Tier"] == "Small") & (df["Nb_Tools"] <= 1)].head(14)
    if qw.empty:
        add_textbox(
            slide, Inches(0.3), Inches(2.5), Inches(9), Inches(0.5),
            "Aucun quick win identifié.", font_size=12, color=PYL_BODY_GREY,
        )
        return

    rows = []
    for _, row in qw.iterrows():
        raw  = str(row["Use Case Description (Long)"])
        desc = raw[:60] + ("..." if len(raw) > 60 else "")
        rows.append([
            str(row["UC_ID"]), str(row["Family_Label"]), str(row["Cluster"]),
            str(row["Tools"]), str(int(row["Score_Total"])),
        ])

    add_pptx_table(
        slide, Inches(0.3), Inches(1.75), Inches(9.4), Inches(3.6),
        ["UC_ID", "Famille", "Cluster", "Outil", "Score"], rows,
    )


def build_recommendations_slide(prs: Presentation, df: pd.DataFrame, page: int = 0) -> None:
    """Slide Recommandations — 4 cartes axes d'industrialisation."""
    slide    = _blank_slide(prs)
    tiers    = df["Complexity_Tier"].value_counts().to_dict()
    it_count = int((df["IT_Flag"] != "").sum())

    add_slide_header(slide, "04", "5 axes d'industrialisation", page=page)

    recos = [
        ("Quick Wins — Déploiement", [
            f"{tiers.get('Small', 0)} UC déployables immédiatement",
            "Bibliothèque de Gems par famille",
            "Cible : 10 pratiques en 3 mois",
        ]),
        ("Accompagner les Medium", [
            f"{tiers.get('Medium', 0)} UC — sprints de 4 semaines",
            "Champion référent par cluster",
            "NotebookLM et App Script comme vecteurs",
        ]),
        (f"Escalader {it_count} cas IT", [
            "Registre de dépendances systèmes",
            "Atelier de priorisation IT",
            "Zero prototype SFDC/SAP sans validation",
        ]),
        ("Track Data Champion", [
            "Parcours formation dédié (SQL, Python, BigQuery)",
            "Pont entre AI builders et data engineers IT",
        ]),
    ]

    CARD_W  = Inches(4.6)
    CARD_H  = Inches(1.85)
    colors  = [PYL_NAVY, PYL_TEAL_BLUE, PYL_DANGER, PYL_YELLOW]
    for i, (title_card, bullets) in enumerate(recos):
        col = i % 2
        row_idx = i // 2
        left = Inches(0.3) + col * (CARD_W + Inches(0.15))
        top  = Inches(1.75) + row_idx * (CARD_H + Inches(0.1))
        add_content_card(slide, left, top, CARD_W, CARD_H, colors[i], title_card, bullets)


def build_closing_slide(prs: Presentation) -> None:
    """Slide de clôture."""
    slide = _blank_slide(prs)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=PYL_NAVY)
    add_textbox(
        slide, Inches(0.5), Inches(2.8), Inches(12), Inches(1.0),
        "Merci de votre attention",
        font_size=36, bold=True, color=PYL_WHITE, align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide, Inches(0.5), Inches(3.9), Inches(12), Inches(0.5),
        "Air Liquide — AI Builders",
        font_size=16, italic=True, color=PYL_YELLOW, align=PP_ALIGN.CENTER,
    )
    add_rect(slide, Inches(3), Inches(4.7), Inches(7.33), Inches(0.05), fill=PYL_YELLOW)


# ── Orchestrateur ─────────────────────────────────────────────────────────────

def main() -> None:
    print(f"Lecture du catalogue : {CATALOG_FILE}")
    if not CATALOG_FILE.exists():
        print(f"Catalogue introuvable : {CATALOG_FILE}")
        sys.exit(1)

    df = pd.read_excel(CATALOG_FILE, sheet_name="Catalogue", header=0)

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    build_cover(prs)
    build_kpi_slide(prs, df, page=2)
    build_complexity_slide(prs, df, page=3)
    build_families_slide(prs, df, page=4)
    build_clusters_slide(prs, df, page=5)
    build_it_slide(prs, df, page=6)
    build_quickwins_slide(prs, df, page=7)
    build_recommendations_slide(prs, df, page=8)
    build_closing_slide(prs)

    OUTPUT_FILE.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_FILE))
    print(f"\nPresentation PowerPoint générée : {OUTPUT_FILE}")


if __name__ == "__main__":
    main()
