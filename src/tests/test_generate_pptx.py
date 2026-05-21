"""
test_generate_pptx.py
Tests unitaires — couverture 100% de src/generate_pptx.py
"""
import sys
import os
import unittest
from pathlib import Path
from unittest.mock import patch, MagicMock

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.path.insert(0, str(Path(__file__).parent.parent))
import generate_pptx as gp


# ── Fixture DataFrame et Presentation ────────────────────────────────────────

def _make_df() -> pd.DataFrame:
    return pd.DataFrame({
        "UC_ID":                        ["UC_0001", "UC_0002", "UC_0003", "UC_0004", "UC_0005"],
        "Cluster":                      ["ClusterA", "ClusterA", "ClusterB", "ClusterC", "ClusterA"],
        "Job Family":                   ["Finance", "Finance", "HR", "Ops", "IT"],
        "Family":                       ["F1", "F7", "F3", "F4", "F1"],
        "Family_Label":                 [
            "Automatisation documentaire",
            "Data Engineering & Reporting",
            "Customer & Sales Intelligence",
            "Monitoring & Maintenance industrielle",
            "Automatisation documentaire",
        ],
        "Complexity_Tier":              ["Small", "Medium", "Large", "Small", "Medium"],
        "Score_Total":                  [5, 9, 13, 6, 8],
        "Score_Integration":            [1, 2, 3, 1, 2],
        "Score_Scope":                  [1, 2, 3, 1, 2],
        "Score_Data":                   [1, 2, 3, 2, 1],
        "Score_AI":                     [1, 1, 2, 1, 1],
        "Score_Economic":               [1, 2, 2, 1, 2],
        "Stage":                        ["POC", "Deployed", "Ideation", "POC", "POC"],
        "Scope of the Use Case":        ["Local", "Country", "Group", "Local", "Local"],
        "Economical Impact":            ["Cost Reduction", "Revenue Growth", "Cost Reduction", "Cost Reduction", "Cost Reduction"],
        "Tools":                        ["Gemini", "Power BI", "AI Studio", "", "Gemini"],
        "Tools_Tags":                   ["Gemini Prompts/Gems", "Power BI", "AI Studio", "", "Gemini Prompts/Gems"],
        "Nb_Tools":                     [1, 1, 1, 0, 1],
        "Max_Tool_Level":               ["L1", "L3", "L3", "L1", "L1"],
        "IT_Flag":                      ["", "⚠️ IT", "⚠️ IT", "", ""],
        "IT_Attention":                 ["", "sfdc", "salesforce", "", ""],
        "Use Case Description (Long)":  [
            "Translate the weekly report to English",
            "Build a bigquery pipeline for reporting",
            "Analyse customer data from sfdc",
            "Monitor equipment sensor data",
            "Draft email using Gemini",
        ],
    })


def _blank_slide(prs: Presentation):
    """Retourne un slide vierge depuis le layout blank (index 6)."""
    return prs.slides.add_slide(prs.slide_layouts[6])


def _new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = gp.SLIDE_W
    prs.slide_height = gp.SLIDE_H
    return prs


# ════════════════════════════════════════════════════════════════════════════
# Helpers bas niveau
# ════════════════════════════════════════════════════════════════════════════

class TestAddRect(unittest.TestCase):
    def test_with_fill_color(self):
        prs  = _new_prs()
        slide = _blank_slide(prs)
        shape = gp.add_rect(slide, Inches(0.5), Inches(0.5), Inches(2), Inches(1),
                             fill=gp.PYL_YELLOW)
        self.assertIsNotNone(shape)

    def test_without_fill_background(self):
        prs  = _new_prs()
        slide = _blank_slide(prs)
        shape = gp.add_rect(slide, Inches(0.5), Inches(0.5), Inches(2), Inches(1),
                             fill=None)
        self.assertIsNotNone(shape)

    def test_no_border(self):
        prs  = _new_prs()
        slide = _blank_slide(prs)
        gp.add_rect(slide, 0, 0, Inches(1), Inches(1), fill=gp.PYL_NAVY, no_border=True)

    def test_with_border(self):
        prs  = _new_prs()
        slide = _blank_slide(prs)
        gp.add_rect(slide, 0, 0, Inches(1), Inches(1), fill=gp.PYL_NAVY, no_border=False)


class TestAddTextbox(unittest.TestCase):
    def test_default_params(self):
        prs  = _new_prs()
        slide = _blank_slide(prs)
        tb = gp.add_textbox(slide, Inches(0.5), Inches(0.5), Inches(3), Inches(0.5), "Hello")
        self.assertIsNotNone(tb)

    def test_with_all_params(self):
        prs  = _new_prs()
        slide = _blank_slide(prs)
        gp.add_textbox(slide, Inches(0.5), Inches(0.5), Inches(3), Inches(0.5),
                       "Styled", font_name="Poppins", font_size=14, bold=True,
                       italic=True, color=gp.PYL_YELLOW, align=PP_ALIGN.CENTER,
                       word_wrap=False)

    def test_without_color(self):
        prs  = _new_prs()
        slide = _blank_slide(prs)
        gp.add_textbox(slide, 0, 0, Inches(2), Inches(0.4), "No color", color=None)


class TestAddLine(unittest.TestCase):
    def test_horizontal_line(self):
        prs  = _new_prs()
        slide = _blank_slide(prs)
        conn = gp.add_line(slide, Inches(0.5), Inches(0.5), Inches(12), Inches(0.5))
        self.assertIsNotNone(conn)

    def test_custom_color_and_width(self):
        prs  = _new_prs()
        slide = _blank_slide(prs)
        gp.add_line(slide, 0, 0, Inches(5), 0, color=gp.PYL_YELLOW, width_pt=2.0)


class TestAddFooter(unittest.TestCase):
    def test_without_page_number(self):
        prs  = _new_prs()
        slide = _blank_slide(prs)
        gp.add_footer(slide, page=0)  # page=0 → pas de numéro

    def test_with_page_number(self):
        prs  = _new_prs()
        slide = _blank_slide(prs)
        gp.add_footer(slide, page=3)  # page>0 → numéro affiché


class TestAddSlideHeader(unittest.TestCase):
    def test_with_subtitle(self):
        prs  = _new_prs()
        slide = _blank_slide(prs)
        gp.add_slide_header(slide, "01", "Titre principal", subtitle="Sous-titre", page=1)

    def test_without_subtitle(self):
        prs  = _new_prs()
        slide = _blank_slide(prs)
        gp.add_slide_header(slide, "02", "Titre sans sous-titre", subtitle="", page=2)


class TestAddKpiCard(unittest.TestCase):
    def test_with_sub(self):
        prs  = _new_prs()
        slide = _blank_slide(prs)
        gp.add_kpi_card(slide, Inches(0.5), Inches(2), Inches(2), Inches(2),
                         "Label", "42", sub="Contexte")

    def test_without_sub(self):
        prs  = _new_prs()
        slide = _blank_slide(prs)
        gp.add_kpi_card(slide, Inches(0.5), Inches(2), Inches(2), Inches(2),
                         "Label", "42")  # sub="" par défaut


class TestAddContentCard(unittest.TestCase):
    def test_with_bullets(self):
        prs  = _new_prs()
        slide = _blank_slide(prs)
        gp.add_content_card(slide, Inches(0.5), Inches(2), Inches(3.5), Inches(4),
                             gp.PYL_NAVY, "Card Title",
                             ["Bullet 1", "Bullet 2", "Bullet 3"])

    def test_empty_bullets(self):
        prs  = _new_prs()
        slide = _blank_slide(prs)
        gp.add_content_card(slide, Inches(0.5), Inches(2), Inches(3.5), Inches(4),
                             gp.PYL_SUCCESS, "Empty card", [])


class TestAddPptxTable(unittest.TestCase):
    def test_table_even_odd_rows(self):
        prs  = _new_prs()
        slide = _blank_slide(prs)
        gp.add_pptx_table(
            slide, Inches(0.5), Inches(2), Inches(12), Inches(3),
            ["Col1", "Col2", "Col3"],
            [["A", "B", "C"], ["D", "E", "F"], ["G", "H", "I"]]  # 3 lignes → even/odd
        )

    def test_single_row(self):
        prs  = _new_prs()
        slide = _blank_slide(prs)
        gp.add_pptx_table(
            slide, Inches(0.5), Inches(2), Inches(12), Inches(1),
            ["X", "Y"],
            [["val1", "val2"]]
        )


# ════════════════════════════════════════════════════════════════════════════
# Builders de slides
# ════════════════════════════════════════════════════════════════════════════

class TestBuildCover(unittest.TestCase):
    def test_adds_one_slide(self):
        prs = _new_prs()
        gp.build_cover(prs)
        self.assertEqual(len(prs.slides), 1)


class TestBuildKpiSlide(unittest.TestCase):
    def test_adds_slide(self):
        prs = _new_prs()
        gp.build_kpi_slide(prs, _make_df(), page=2)
        self.assertEqual(len(prs.slides), 1)


class TestBuildComplexitySlide(unittest.TestCase):
    def test_adds_slide(self):
        prs = _new_prs()
        gp.build_complexity_slide(prs, _make_df(), page=3)
        self.assertEqual(len(prs.slides), 1)


class TestBuildFamiliesSlide(unittest.TestCase):
    def test_adds_slide(self):
        prs = _new_prs()
        gp.build_families_slide(prs, _make_df(), page=4)
        self.assertEqual(len(prs.slides), 1)


class TestBuildClustersSlide(unittest.TestCase):
    def test_adds_slide(self):
        prs = _new_prs()
        gp.build_clusters_slide(prs, _make_df(), page=5)
        self.assertEqual(len(prs.slides), 1)

    def test_handles_missing_tier_columns(self):
        """Si Small/Medium/Large est absent du groupby, ne pas planter."""
        df = _make_df()
        df["Complexity_Tier"] = "Medium"  # Plus de Small ni Large
        prs = _new_prs()
        gp.build_clusters_slide(prs, df, page=5)


class TestBuildItSlide(unittest.TestCase):
    def test_with_it_flags(self):
        prs = _new_prs()
        gp.build_it_slide(prs, _make_df(), page=6)
        self.assertEqual(len(prs.slides), 1)

    def test_with_no_it_flags(self):
        df = _make_df()
        df["IT_Flag"] = ""
        prs = _new_prs()
        gp.build_it_slide(prs, df, page=6)

    def test_long_description_truncated(self):
        """Couvre la branche 'desc += ...' pour les descriptions IT > 50 chars."""
        df = _make_df()
        df["Use Case Description (Long)"] = "X" * 60  # > 50 chars
        df["IT_Flag"] = "\u26a0\ufe0f IT"
        prs = _new_prs()
        gp.build_it_slide(prs, df, page=6)


class TestBuildQuickwinsSlide(unittest.TestCase):
    def test_with_quick_wins(self):
        prs = _new_prs()
        gp.build_quickwins_slide(prs, _make_df(), page=7)
        self.assertEqual(len(prs.slides), 1)

    def test_with_no_quick_wins(self):
        df = _make_df()
        df["Complexity_Tier"] = "Large"  # Aucun Small
        prs = _new_prs()
        gp.build_quickwins_slide(prs, df, page=7)

    def test_long_description_truncated(self):
        """Couvre la branche 'desc += ...' pour les descriptions Quick Wins > 60 chars."""
        df = _make_df()
        df["Use Case Description (Long)"] = "Y" * 70  # > 60 chars
        df["Complexity_Tier"] = "Small"
        df["Nb_Tools"] = 1
        prs = _new_prs()
        gp.build_quickwins_slide(prs, df, page=7)


class TestBuildRecommendationsSlide(unittest.TestCase):
    def test_adds_slide(self):
        prs = _new_prs()
        gp.build_recommendations_slide(prs, _make_df(), page=8)
        self.assertEqual(len(prs.slides), 1)


class TestBuildClosingSlide(unittest.TestCase):
    def test_adds_slide(self):
        prs = _new_prs()
        gp.build_closing_slide(prs)
        self.assertEqual(len(prs.slides), 1)


# ════════════════════════════════════════════════════════════════════════════
# Pipeline complet
# ════════════════════════════════════════════════════════════════════════════

class TestFullPipeline(unittest.TestCase):
    def test_all_slides_generated(self):
        prs = _new_prs()
        df  = _make_df()
        gp.build_cover(prs)
        gp.build_kpi_slide(prs, df, page=2)
        gp.build_complexity_slide(prs, df, page=3)
        gp.build_families_slide(prs, df, page=4)
        gp.build_clusters_slide(prs, df, page=5)
        gp.build_it_slide(prs, df, page=6)
        gp.build_quickwins_slide(prs, df, page=7)
        gp.build_recommendations_slide(prs, df, page=8)
        gp.build_closing_slide(prs)
        self.assertEqual(len(prs.slides), 9)


# ════════════════════════════════════════════════════════════════════════════
# main() — filesystem mocking
# ════════════════════════════════════════════════════════════════════════════

class TestMain(unittest.TestCase):
    def test_exits_when_catalog_missing(self):
        fake_source = MagicMock(spec=Path)
        fake_source.exists.return_value = False

        with patch.object(gp, "CATALOG_FILE", fake_source):
            with self.assertRaises(SystemExit) as ctx:
                gp.main()
        self.assertEqual(ctx.exception.code, 1)

    def test_success_generates_pptx(self):
        import tempfile
        df = _make_df()
        fake_source = MagicMock(spec=Path)
        fake_source.exists.return_value = True

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            tmp_path = Path(tmp.name)

        try:
            with patch.object(gp, "CATALOG_FILE", fake_source), \
                 patch.object(gp, "OUTPUT_FILE", tmp_path), \
                 patch("generate_pptx.pd.read_excel", return_value=df):
                gp.main()
            self.assertTrue(tmp_path.exists())
            self.assertGreater(tmp_path.stat().st_size, 5000)
        finally:
            if tmp_path.exists():
                os.unlink(tmp_path)


if __name__ == "__main__":
    unittest.main(verbosity=2)
