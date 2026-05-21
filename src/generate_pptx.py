#!/usr/bin/env python3
"""
generate_pptx.py
────────────────
Genere la presentation PowerPoint (.pptx) Air Liquide — AI Champions
avec la charte graphique officielle Pyl.Tech.

Style : Presentation de consulting — cabinet Tier-1 (contexte mission,
        methodologie, framework d'analyse, architecture, recommandations).

Usage : python3 src/generate_pptx.py
Output: docs/presentation_ai_champions.pptx
"""
import sys
from pathlib import Path
from datetime import date

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    import pandas as pd
except ImportError as e:  # pragma: no cover
    print(f"Dependance manquante : {e}\n   pip install python-pptx pandas openpyxl")  # noqa
    sys.exit(1)

# ── Chemins ───────────────────────────────────────────────────────────────────
PROJECT_DIR = Path(__file__).parent.parent
CATALOG_FILE = PROJECT_DIR / "output" / "use_cases_catalog.xlsx"
OUTPUT_FILE = PROJECT_DIR / "output" / "presentation_ai_champions.pptx"

# ── Palette Pyl.Tech officielle ───────────────────────────────────────────────
PYL_NAVY_DARK = RGBColor(0x0B, 0x13, 0x2B)
PYL_NAVY = RGBColor(0x0D, 0x21, 0x49)
PYL_YELLOW = RGBColor(0xF4, 0xBF, 0x46)
PYL_TEAL_BLUE = RGBColor(0x20, 0x8A, 0xAE)
PYL_TEAL = RGBColor(0x5B, 0xC0, 0xBE)
PYL_BODY_GREY = RGBColor(0x4F, 0x4F, 0x4F)
PYL_GREY_BG = RGBColor(0xEE, 0xEE, 0xEE)
PYL_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PYL_SUCCESS = RGBColor(0x13, 0x86, 0x36)
PYL_DANGER = RGBColor(0xC9, 0x14, 0x32)

YEAR = date.today().year
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── Helpers de bas niveau ─────────────────────────────────────────────────────


def add_rect(slide, left, top, width, height, fill: RGBColor = None, no_border: bool = True):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if no_border:
        shape.line.fill.background()
    return shape


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    font_name: str = "Poppins",
    font_size: int = 12,
    bold: bool = False,
    italic: bool = False,
    color: RGBColor = None,
    align=PP_ALIGN.LEFT,
    word_wrap: bool = True,
):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = word_wrap
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return txBox


def add_multiline_textbox(
    slide,
    left,
    top,
    width,
    height,
    lines: list[tuple[str, int, bool, RGBColor]],
    align=PP_ALIGN.LEFT,
):
    """Textbox with multiple paragraphs. Each line = (text, size, bold, color)."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    tf = txBox.text_frame

    for i, (text, size, bold, color) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name = "Poppins"
        run.font.size = Pt(size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color

    return txBox


def add_line(slide, x1, y1, x2, y2, color: RGBColor = PYL_NAVY_DARK, width_pt: float = 0.75):
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = Pt(width_pt)
    return connector


def add_slide_header(
    slide, section_num: str, title: str, subtitle: str = "", page: int = 0
) -> None:
    """Header Pyl.Tech : cartouche jaune + titre + traits + sous-titre."""
    MARGIN_L = Inches(0.5)
    MARGIN_R = Inches(12.83)

    add_line(slide, MARGIN_L, Inches(0.45), MARGIN_R, Inches(0.45))

    badge = add_rect(
        slide, MARGIN_L, Inches(0.5), Inches(0.55), Inches(0.55), fill=PYL_YELLOW
    )
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = section_num
    run.font.name = "Poppins"
    run.font.bold = True
    run.font.size = Pt(16)
    run.font.color.rgb = PYL_WHITE

    title_box = slide.shapes.add_textbox(
        Inches(1.2), Inches(0.48), Inches(11.0), Inches(0.65)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Poppins"
    run.font.bold = True
    run.font.size = Pt(26)
    run.font.color.rgb = PYL_NAVY_DARK

    if subtitle:
        sub_box = slide.shapes.add_textbox(
            Inches(1.2), Inches(1.15), Inches(11.0), Inches(0.4)
        )
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = subtitle
        run.font.name = "Poppins"
        run.font.size = Pt(13)
        run.font.color.rgb = PYL_BODY_GREY

    bottom_y = Inches(1.6) if subtitle else Inches(1.25)
    add_line(slide, MARGIN_L, bottom_y, MARGIN_R, bottom_y)

    add_textbox(
        slide,
        Inches(11.8),
        Inches(0.15),
        Inches(1.4),
        Inches(0.35),
        "Pyl.Tech",
        font_size=11,
        bold=True,
        color=PYL_NAVY_DARK,
        align=PP_ALIGN.RIGHT,
    )

    add_footer(slide, page)


def add_footer(slide, page: int = 0) -> None:
    add_textbox(
        slide,
        Inches(0.3),
        Inches(7.15),
        Inches(5),
        Inches(0.25),
        f"(c) Copyright {YEAR} Pyl.Tech  |  Air Liquide AI Champions",
        font_size=9,
        color=PYL_TEAL_BLUE,
    )
    if page:
        add_textbox(
            slide,
            Inches(12.5),
            Inches(7.15),
            Inches(0.7),
            Inches(0.25),
            str(page),
            font_size=9,
            color=PYL_TEAL_BLUE,
            align=PP_ALIGN.RIGHT,
        )


def add_kpi_card(
    slide,
    left,
    top,
    width,
    height,
    label: str,
    value: str,
    sub: str = "",
    header_color: RGBColor = PYL_NAVY,
    value_color: RGBColor = PYL_TEAL_BLUE,
) -> None:
    card = add_rect(slide, left, top, width, height, fill=PYL_WHITE)
    card.line.color.rgb = PYL_GREY_BG
    card.line.width = Pt(0.5)

    add_rect(slide, left, top, width, Inches(0.07), fill=header_color)

    v_box = slide.shapes.add_textbox(
        left + Inches(0.1), top + Inches(0.15), width - Inches(0.2), Inches(0.9)
    )
    v_box.text_frame.word_wrap = False
    p = v_box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = value
    run.font.name = "Poppins"
    run.font.bold = True
    run.font.size = Pt(36)
    run.font.color.rgb = value_color

    l_box = slide.shapes.add_textbox(
        left + Inches(0.1), top + Inches(1.0), width - Inches(0.2), Inches(0.5)
    )
    l_box.text_frame.word_wrap = True
    p = l_box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.name = "Poppins"
    run.font.size = Pt(10)
    run.font.color.rgb = PYL_BODY_GREY

    if sub:
        s_box = slide.shapes.add_textbox(
            left + Inches(0.1),
            top + Inches(1.45),
            width - Inches(0.2),
            Inches(0.35),
        )
        s_box.text_frame.word_wrap = True
        p = s_box.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = sub
        run.font.name = "Poppins"
        run.font.size = Pt(9)
        run.font.italic = True
        run.font.color.rgb = PYL_BODY_GREY


def add_content_card(
    slide,
    left,
    top,
    width,
    height,
    header_color: RGBColor,
    card_title: str,
    bullets: list[str],
) -> None:
    card_bg = add_rect(slide, left, top, width, height, fill=PYL_WHITE)
    card_bg.line.color.rgb = PYL_GREY_BG
    card_bg.line.width = Pt(0.5)

    add_rect(slide, left, top, width, Inches(0.45), fill=header_color)

    t_box = slide.shapes.add_textbox(
        left + Inches(0.15), top + Inches(0.06), width - Inches(0.3), Inches(0.35)
    )
    p = t_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = card_title
    run.font.name = "Poppins"
    run.font.bold = True
    run.font.size = Pt(13)
    run.font.color.rgb = PYL_WHITE

    cur_top = top + Inches(0.55)
    for bullet in bullets:
        b_box = slide.shapes.add_textbox(
            left + Inches(0.15), cur_top, width - Inches(0.3), Inches(0.35)
        )
        b_box.text_frame.word_wrap = True
        p = b_box.text_frame.paragraphs[0]
        sq = p.add_run()
        sq.text = "  "
        sq.font.name = "Poppins"
        sq.font.bold = True
        sq.font.size = Pt(11)
        sq.font.color.rgb = PYL_YELLOW
        txt = p.add_run()
        txt.text = bullet
        txt.font.name = "Poppins"
        txt.font.size = Pt(11)
        txt.font.color.rgb = PYL_BODY_GREY
        cur_top += Inches(0.42)


def add_pptx_table(
    slide, left, top, width, height, headers: list[str], rows: list[list[str]]
) -> None:
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table

    col_w = width // n_cols
    for i in range(n_cols):
        tbl.columns[i].width = col_w

    for i, h in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PYL_NAVY_DARK
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = h
        run.font.name = "Poppins"
        run.font.bold = True
        run.font.size = Pt(10)
        run.font.color.rgb = PYL_WHITE

    for r_idx, row_data in enumerate(rows):
        bg = PYL_NAVY if r_idx % 2 == 0 else PYL_WHITE
        fg = PYL_WHITE if r_idx % 2 == 0 else PYL_BODY_GREY
        for c_idx, val in enumerate(row_data):
            cell = tbl.cell(r_idx + 1, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if c_idx > 0 else PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(val)
            run.font.name = "Poppins"
            run.font.size = Pt(10)
            run.font.color.rgb = fg


# ── Construction des slides ───────────────────────────────────────────────────


def build_cover(prs: Presentation) -> None:
    """Slide de couverture — fond jaune integral."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=PYL_YELLOW)

    add_textbox(
        slide,
        Inches(0.5),
        Inches(0.3),
        Inches(2),
        Inches(0.45),
        "Pyl.Tech",
        font_size=18,
        bold=True,
        color=PYL_NAVY_DARK,
    )

    t_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(2.0), Inches(10.3), Inches(1.6)
    )
    t_box.text_frame.word_wrap = True
    p = t_box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Air Liquide — AI Champions"
    run.font.name = "Poppins"
    run.font.bold = True
    run.font.size = Pt(44)
    run.font.color.rgb = PYL_NAVY_DARK

    st_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(3.6), Inches(10.3), Inches(0.7)
    )
    st_box.text_frame.word_wrap = True
    p = st_box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Analyse, Classification et Recommandations\ndu Portefeuille Use Cases IA"
    run.font.name = "Poppins"
    run.font.italic = True
    run.font.size = Pt(20)
    run.font.color.rgb = PYL_NAVY_DARK

    add_rect(
        slide,
        Inches(3),
        Inches(4.85),
        Inches(7.3),
        Inches(0.03),
        fill=PYL_NAVY_DARK,
    )

    add_textbox(
        slide,
        Inches(1.5),
        Inches(5.1),
        Inches(10.3),
        Inches(0.4),
        f"Produit par Pyl.Tech  |  {date.today().strftime('%B %Y')}",
        font_size=13,
        color=PYL_NAVY,
        align=PP_ALIGN.CENTER,
    )

    add_textbox(
        slide,
        Inches(1.5),
        Inches(5.6),
        Inches(10.3),
        Inches(0.4),
        "Document confidentiel — Usage interne Air Liquide",
        font_size=11,
        italic=True,
        color=PYL_BODY_GREY,
        align=PP_ALIGN.CENTER,
    )

    add_textbox(
        slide,
        Inches(0.3),
        Inches(7.1),
        Inches(7),
        Inches(0.25),
        f"(c) Copyright {YEAR} Pyl.Tech  |  Confidentiel",
        font_size=9,
        color=PYL_NAVY_DARK,
    )


def build_sommaire(prs: Presentation, page: int) -> None:
    """Slide Sommaire — plan du document."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(
        slide, "", "Sommaire", "Plan de la presentation", page
    )

    sections = [
        ("01", "Contexte de la mission"),
        ("02", "Méthodologie d'analyse"),
        ("03", "Chiffres clés du portefeuille"),
        ("04", "Répartition par complexité"),
        ("05", "Familles fonctionnelles"),
        ("06", "Répartition par Cluster"),
        ("07", "Points d'attention IT"),
        ("08", "Architectures de reference"),
        ("09", "Top Quick Wins"),
        ("10", "Recommandations et Roadmap"),
        ("11", "Bonnes pratiques AI Champions"),
    ]

    for i, (num, title) in enumerate(sections):
        y = Inches(1.85) + i * Inches(0.46)

        badge = add_rect(
            slide, Inches(1.5), y, Inches(0.55), Inches(0.38), fill=PYL_YELLOW
        )
        bt = badge.text_frame
        bp = bt.paragraphs[0]
        bp.alignment = PP_ALIGN.CENTER
        br = bp.add_run()
        br.text = num
        br.font.name = "Poppins"
        br.font.bold = True
        br.font.size = Pt(11)
        br.font.color.rgb = PYL_NAVY_DARK

        add_textbox(
            slide,
            Inches(2.2),
            y,
            Inches(9),
            Inches(0.38),
            title,
            font_size=14,
            color=PYL_NAVY_DARK,
        )


def build_context_slide(prs: Presentation, df: pd.DataFrame, page: int) -> None:
    """Slide Contexte de la mission."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(
        slide,
        "01",
        "Contexte de la mission",
        "Cadrage du projet et objectifs de l'analyse",
        page,
    )

    total = len(df)

    # Contexte text
    add_textbox(
        slide,
        Inches(0.5),
        Inches(1.8),
        Inches(6.0),
        Inches(1.2),
        f"Air Liquide a mobilise ses AI Champions — profils metiers, non-IT — "
        f"pour explorer l'IA dans leurs activites. {total} use cases "
        f"ont ete collectes sur l'ensemble des clusters du groupe.",
        font_size=13,
        color=PYL_BODY_GREY,
    )

    # 4 objectifs en cards
    objectives = [
        ("Comprendre", "Regroupement thematique en familles fonctionnelles transverses"),
        ("Qualifier", "Scoring multi-dimensionnel : Small / Medium / Large"),
        ("Projeter", "Architecture cible Google-first par niveau de complexité"),
        ("Guider", "Recommandations et bonnes pratiques pour les champions"),
    ]

    CARD_W = Inches(2.9)
    CARD_H = Inches(2.5)
    for i, (obj_title, obj_desc) in enumerate(objectives):
        col = i % 2
        row_idx = i // 2
        left = Inches(0.5) + col * (CARD_W + Inches(0.2))
        top = Inches(3.3) + row_idx * (CARD_H + Inches(0.15))

        colors = [PYL_NAVY, PYL_TEAL_BLUE, PYL_TEAL, PYL_YELLOW]
        add_content_card(
            slide, left, top, CARD_W, CARD_H, colors[i], obj_title, [obj_desc]
        )

    # Callout droite
    add_rect(
        slide,
        Inches(6.7),
        Inches(3.3),
        Inches(6.0),
        Inches(3.0),
        fill=PYL_GREY_BG,
    )
    add_textbox(
        slide,
        Inches(6.9),
        Inches(3.5),
        Inches(5.6),
        Inches(0.4),
        "Périmètre de l'analyse",
        font_size=16,
        bold=True,
        color=PYL_NAVY_DARK,
    )

    perim_items = [
        f"{total} use cases analyses",
        f"{df['Cluster'].nunique()} clusters géographiques/organisationnels",
        f"{df['Job Family'].nunique()} familles metier",
        "11 outils distincts identifies",
        "Source : Advanced AI Champions - Action Monitoring.xlsx",
    ]
    for j, item in enumerate(perim_items):
        add_textbox(
            slide,
            Inches(6.9),
            Inches(4.1) + j * Inches(0.4),
            Inches(5.6),
            Inches(0.35),
            f"  {item}",
            font_size=12,
            color=PYL_BODY_GREY,
        )


def build_methodology_slide(prs: Presentation, page: int) -> None:
    """Slide Méthodologie — 6 phases."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(
        slide,
        "02",
        "Méthodologie d'analyse",
        "Approche en 6 phases pour qualifier le portefeuille AI Champions",
        page,
    )

    phases = [
        ("Phase 1", "Audit & Nettoyage", "Traitement des anomalies, normalisation, UC_ID", PYL_NAVY),
        ("Phase 2", "Classification", "Analyse semantique en 7 familles fonctionnelles", PYL_TEAL_BLUE),
        ("Phase 3", "Scoring", "5 dimensions x 3 niveaux = score 5-15", PYL_TEAL),
        ("Phase 4", "Detection IT", "Identification des dependances systèmes", PYL_DANGER),
        ("Phase 5", "Architecture", "Stack Google-first par tier de complexité", PYL_YELLOW),
        ("Phase 6", "Recommandations", "Bonnes pratiques et feuille de route", PYL_SUCCESS),
    ]

    # 3 cards per row, 2 rows
    CARD_W = Inches(3.9)
    CARD_H = Inches(2.3)

    for i, (phase_num, phase_title, phase_desc, color) in enumerate(phases):
        col = i % 3
        row_idx = i // 3
        left = Inches(0.5) + col * (CARD_W + Inches(0.15))
        top = Inches(1.85) + row_idx * (CARD_H + Inches(0.15))

        add_content_card(
            slide,
            left,
            top,
            CARD_W,
            CARD_H,
            color,
            f"{phase_num} — {phase_title}",
            [phase_desc],
        )

    # Fleche de processus en bas
    add_textbox(
        slide,
        Inches(0.5),
        Inches(6.6),
        Inches(12.3),
        Inches(0.4),
        "Audit  ->  Classification  ->  Scoring  ->  Detection IT  ->  Architecture  ->  Recommandations",
        font_size=12,
        bold=True,
        color=PYL_TEAL_BLUE,
        align=PP_ALIGN.CENTER,
    )


def build_framework_slide(prs: Presentation, page: int) -> None:
    """Slide Framework de scoring — tableau des 5 dimensions."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(
        slide,
        "02",
        "Cadre d'analyse : scoring multi-dimensionnel",
        "5 dimensions x 3 niveaux = score 5 a 15 points",
        page,
    )

    add_pptx_table(
        slide,
        Inches(0.5),
        Inches(1.85),
        Inches(12.3),
        Inches(3.8),
        ["Dimension", "1 pt (Faible)", "2 pts (Moyen)", "3 pts (Élevé)"],
        [
            [
                "D1 Integration technique",
                "1 outil, no-code (L1/L2)",
                "2-3 outils, semi-code (L3)",
                "4+ outils, code custom (L4)",
            ],
            [
                "D2 Périmètre organisationnel",
                "Equipe locale",
                "Country / Cluster",
                "Group (global)",
            ],
            [
                "D3 Complexite data",
                "Donnees statiques",
                "Donnees connectees (SFDC, BI)",
                "Temps reel / industrielles",
            ],
            [
                "D4 Maturite IA",
                "Prompting / Gem",
                "API Gemini + RAG",
                "Agent / ML / fine-tuning",
            ],
            [
                "D5 Impact économique",
                "Non évalué / Productivite",
                "Cost Reduction",
                "Revenue Growth",
            ],
        ],
    )

    # Tiers en 3 mini-cards
    tiers_info = [
        ("Small (5-7)", "Quick Win", "< 2 sem.", "Champion seul", PYL_SUCCESS),
        ("Medium (8-11)", "Structurant", "4-8 sem.", "Champion + IT", PYL_TEAL_BLUE),
        ("Large (12-15)", "Stratégique", "3-12 mois", "Equipe IT", PYL_DANGER),
    ]

    CARD_W = Inches(3.9)
    for i, (tier, label, ttv, profile, color) in enumerate(tiers_info):
        left = Inches(0.5) + i * (CARD_W + Inches(0.27))
        top = Inches(6.0)

        badge = add_rect(slide, left, top, CARD_W, Inches(0.6), fill=color)
        bt = badge.text_frame
        bt.word_wrap = True
        bp = bt.paragraphs[0]
        bp.alignment = PP_ALIGN.CENTER
        br = bp.add_run()
        br.text = f"{tier}  |  {label}  |  {ttv}  |  {profile}"
        br.font.name = "Poppins"
        br.font.bold = True
        br.font.size = Pt(10)
        br.font.color.rgb = PYL_WHITE


def build_kpi_slide(prs: Presentation, df: pd.DataFrame, page: int) -> None:
    """Slide chiffres clés — 6 KPI cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(
        slide,
        "03",
        "Chiffres clés du portefeuille",
        f"{len(df)} use cases analyses sur l'ensemble des clusters Air Liquide",
        page,
    )

    total = len(df)
    tiers = df["Complexity_Tier"].value_counts().to_dict()
    it_count = (df["IT_Flag"] != "").sum()
    clusters = df["Cluster"].nunique()

    CARD_W = Inches(2.05)
    CARD_H = Inches(1.95)
    Y_TOP = Inches(1.85)
    GAP = Inches(0.1)

    cards = [
        ("Use Cases analyses", str(total), "Portefeuille complet", PYL_NAVY, PYL_YELLOW),
        ("Clusters couverts", str(clusters), "Périmètre organisationnel", PYL_NAVY, PYL_TEAL_BLUE),
        ("Quick Wins (Small)", str(tiers.get("Small", 0)), "< 2 semaines", PYL_SUCCESS, PYL_SUCCESS),
        ("Structurants (Medium)", str(tiers.get("Medium", 0)), "4-8 semaines", PYL_TEAL, PYL_TEAL),
        ("Stratégiques (Large)", str(tiers.get("Large", 0)), "3-12 mois", PYL_DANGER, PYL_DANGER),
        ("Points d'attention IT", str(it_count), "Escalade requise", PYL_DANGER, PYL_DANGER),
    ]

    for i, (label, value, sub, hdr_color, val_color) in enumerate(cards):
        left = Inches(0.5) + i * (CARD_W + GAP)
        add_kpi_card(
            slide,
            left,
            Y_TOP,
            CARD_W,
            CARD_H,
            label,
            value,
            sub,
            header_color=hdr_color,
            value_color=val_color,
        )

    pct_accessible = (tiers.get("Small", 0) + tiers.get("Medium", 0)) * 100 // total
    add_textbox(
        slide,
        Inches(0.5),
        Inches(4.1),
        Inches(12.3),
        Inches(0.4),
        f"{pct_accessible}% des use cases sont deployables sans ressources IT lourdes (Small ou Medium)",
        font_size=13,
        bold=True,
        color=PYL_NAVY_DARK,
        align=PP_ALIGN.CENTER,
    )

    callout = add_rect(
        slide,
        Inches(0.5),
        Inches(4.6),
        Inches(12.3),
        Inches(1.5),
        fill=PYL_YELLOW,
    )
    ct = callout.text_frame
    ct.word_wrap = True
    p = ct.paragraphs[0]
    run = p.add_run()
    run.text = (
        f"Constat clé : la majorité du portefeuille ({tiers.get('Small', 0) + tiers.get('Medium', 0)} UC) "
        f"représente des opportunités de déploiement rapide avec Google Workspace. "
        f"Seuls {tiers.get('Large', 0)} projets nécessitent un engagement long terme avec l'IT. "
        f"{it_count} UC ont des dependances systèmes qui requièrent une escalade obligatoire."
    )
    run.font.name = "Poppins"
    run.font.size = Pt(12)
    run.font.color.rgb = PYL_NAVY_DARK


def build_complexity_slide(prs: Presentation, df: pd.DataFrame, page: int) -> None:
    """Slide Répartition par complexité — 3 cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(
        slide,
        "04",
        "Répartition par complexité",
        "Scoring 5 dimensions (5-15 pts) -> Small / Medium / Large",
        page,
    )

    tiers = df["Complexity_Tier"].value_counts().to_dict()
    total = len(df)

    tier_data = [
        ("Small", tiers.get("Small", 0), "Quick Win", "< 2 sem.", "1 champion, no-code", PYL_SUCCESS),
        ("Medium", tiers.get("Medium", 0), "Use Case Structurant", "4-8 sem.", "Champion + IT local", PYL_TEAL_BLUE),
        ("Large", tiers.get("Large", 0), "Projet Stratégique", "3-12 mois", "Equipe IT + Champion", PYL_DANGER),
    ]

    CARD_W = Inches(3.9)
    CARD_H = Inches(4.5)
    Y_TOP = Inches(1.8)

    for i, (tier, count, label, ttv, profile, color) in enumerate(tier_data):
        left = Inches(0.5) + i * (CARD_W + Inches(0.27))
        pct = count * 100 // total if total > 0 else 0
        add_content_card(
            slide,
            left,
            Y_TOP,
            CARD_W,
            CARD_H,
            color,
            f"{tier} — {count} use cases ({pct}%)",
            [
                f"Label : {label}",
                f"Time-to-value : {ttv}",
                f"Profil : {profile}",
                f"Part du portefeuille : {pct}%",
            ],
        )

        add_textbox(
            slide,
            left,
            Y_TOP + Inches(0.55),
            CARD_W,
            Inches(0.8),
            str(count),
            font_size=48,
            bold=True,
            color=PYL_WHITE,
            align=PP_ALIGN.CENTER,
        )


def build_families_slide(prs: Presentation, df: pd.DataFrame, page: int) -> None:
    """Slide Familles fonctionnelles."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(
        slide,
        "05",
        "Familles fonctionnelles",
        "7 familles identifiées — répartition et volume par tier",
        page,
    )

    families = df["Family_Label"].value_counts().reset_index()
    families.columns = ["Family_Label", "Total"]
    families = families.sort_values("Total", ascending=False)

    rows = []
    for _, row in families.iterrows():
        fam = row["Family_Label"]
        subset = (
            df[df["Family_Label"] == fam]["Complexity_Tier"]
            .value_counts()
            .to_dict()
        )
        code = df[df["Family_Label"] == fam]["Family"].iloc[0]
        rows.append(
            [
                f"{code} — {fam}",
                str(subset.get("Small", 0)),
                str(subset.get("Medium", 0)),
                str(subset.get("Large", 0)),
                str(row["Total"]),
            ]
        )

    add_pptx_table(
        slide,
        Inches(0.5),
        Inches(1.9),
        Inches(12.3),
        Inches(4.8),
        ["Famille fonctionnelle", "Small", "Medium", "Large", "Total"],
        rows,
    )


def build_clusters_slide(prs: Presentation, df: pd.DataFrame, page: int) -> None:
    """Slide Top clusters."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(
        slide,
        "06",
        "Répartition par Cluster",
        "Top 12 clusters par volume de use cases",
        page,
    )

    cluster_agg = (
        df.groupby("Cluster")["Complexity_Tier"]
        .value_counts()
        .unstack(fill_value=0)
        .reset_index()
    )
    cluster_agg["Total"] = (
        cluster_agg.get("Small", 0)
        + cluster_agg.get("Medium", 0)
        + cluster_agg.get("Large", 0)
    )
    cluster_agg = cluster_agg.sort_values("Total", ascending=False).head(12)

    rows = []
    for _, row in cluster_agg.iterrows():
        rows.append(
            [
                str(row["Cluster"]),
                str(int(row.get("Small", 0))),
                str(int(row.get("Medium", 0))),
                str(int(row.get("Large", 0))),
                str(int(row["Total"])),
            ]
        )

    add_pptx_table(
        slide,
        Inches(0.5),
        Inches(1.9),
        Inches(12.3),
        Inches(4.8),
        ["Cluster", "Small", "Medium", "Large", "Total"],
        rows,
    )


def build_it_slide(prs: Presentation, df: pd.DataFrame, page: int) -> None:
    """Slide Points d'attention IT."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    it_count = (df["IT_Flag"] != "").sum()
    add_slide_header(
        slide,
        "07",
        f"{it_count} points d'attention IT",
        "Use cases necessitant un accompagnement technique dédié",
        page,
    )

    alert = add_rect(
        slide,
        Inches(0.5),
        Inches(1.75),
        Inches(12.3),
        Inches(0.65),
        fill=PYL_DANGER,
    )
    alert.line.fill.background()
    at = alert.text_frame
    at.word_wrap = True
    p = at.paragraphs[0]
    run = p.add_run()
    run.text = (
        "Ces use cases impliquent SFDC, SAP, DCS, SCADA, BigQuery, API d'entreprise ou profil Large. "
        "Escalade IT obligatoire AVANT tout prototype."
    )
    run.font.name = "Poppins"
    run.font.bold = True
    run.font.size = Pt(12)
    run.font.color.rgb = PYL_WHITE

    it_df = df[df["IT_Flag"] != ""]
    it_by_fam = (
        it_df.groupby("Family_Label")
        .size()
        .reset_index(name="count")
        .sort_values("count", ascending=False)
    )

    rows = [[r["Family_Label"], str(r["count"])] for _, r in it_by_fam.iterrows()]

    add_pptx_table(
        slide,
        Inches(0.5),
        Inches(2.6),
        Inches(6),
        Inches(4.1),
        ["Famille fonctionnelle", "Nb use cases IT"],
        rows,
    )

    top_it = it_df.nlargest(8, "Score_Total")
    right_rows = []
    for _, row in top_it.iterrows():
        desc = str(row["Use Case Description (Long)"])[:50]
        if len(str(row["Use Case Description (Long)"])) > 50:
            desc += "..."
        right_rows.append(
            [
                str(row["UC_ID"]),
                str(row["Cluster"]),
                str(row["Complexity_Tier"]),
                desc,
            ]
        )

    add_pptx_table(
        slide,
        Inches(6.8),
        Inches(2.6),
        Inches(6.0),
        Inches(4.1),
        ["UC_ID", "Cluster", "Tier", "Description"],
        right_rows,
    )


def build_architecture_slide(prs: Presentation, page: int) -> None:
    """Slide Architecture cible — 3 tiers."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(
        slide,
        "08",
        "Architectures de reference",
        "Stack cible Google-first par niveau de complexité",
        page,
    )

    arch_data = [
        (
            PYL_SUCCESS,
            "Small — Prompting & Automation",
            [
                "Outils : Gemini Gems, NotebookLM",
                "Data : Google Drive / Sheets (statique)",
                "Output : Docs / Gmail / Chat",
                "Competence : Prompt engineering",
                "IT : Aucune intervention requise",
            ],
        ),
        (
            PYL_TEAL_BLUE,
            "Medium — App & Orchestration",
            [
                "Outils : App Script, AppSheet, AI Studio",
                "Data : Sheets, AppSheet (structurees)",
                "Output : AppSheet App, Slides/Docs",
                "Competence : Low-code + API basiques",
                "IT : Support ponctuel si SFDC/BI",
            ],
        ),
        (
            PYL_DANGER,
            "Large — Platform & Agent",
            [
                "Sources : SFDC, AVEVA, DCS, SAP, Fabric",
                "Ingestion : Python, BigQuery, DataStudio",
                "AI : Vertex AI + RAG + Agents Gemini",
                "Backend : Cloud Run, Advance Coding",
                "IT : Equipe projet dédiée obligatoire",
            ],
        ),
    ]

    CARD_W = Inches(3.9)
    CARD_H = Inches(4.5)

    for i, (color, title, bullets) in enumerate(arch_data):
        left = Inches(0.5) + i * (CARD_W + Inches(0.27))
        add_content_card(
            slide, left, Inches(1.8), CARD_W, CARD_H, color, title, bullets
        )


def build_quickwins_slide(prs: Presentation, df: pd.DataFrame, page: int) -> None:
    """Slide Top Quick Wins."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(
        slide,
        "09",
        "Top Quick Wins",
        "Use cases Small (1 outil) a deployer en priorite",
        page,
    )

    qw = df[(df["Complexity_Tier"] == "Small") & (df["Nb_Tools"] <= 1)].head(10)

    rows = []
    for _, row in qw.iterrows():
        desc = str(row["Use Case Description (Long)"])[:60]
        if len(str(row["Use Case Description (Long)"])) > 60:
            desc += "..."
        rows.append(
            [
                str(row["UC_ID"]),
                str(row["Family_Label"]),
                str(row["Cluster"]),
                str(row["Tools"]),
                str(int(row["Score_Total"])),
                desc,
            ]
        )

    add_pptx_table(
        slide,
        Inches(0.5),
        Inches(1.9),
        Inches(12.3),
        Inches(4.8),
        ["UC_ID", "Famille", "Cluster", "Outil", "Score", "Description"],
        rows,
    )


def build_recommendations_slide(
    prs: Presentation, df: pd.DataFrame, page: int
) -> None:
    """Slide Recommandations — 4 cards + roadmap."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(
        slide,
        "10",
        "Recommandations et Roadmap",
        "5 axes prioritaires pour industrialiser les use cases",
        page,
    )

    tiers = df["Complexity_Tier"].value_counts().to_dict()
    it_count = (df["IT_Flag"] != "").sum()

    recs = [
        (
            PYL_SUCCESS,
            "Industrialiser les Quick Wins",
            [
                f"{tiers.get('Small', 0)} UC deployables immediatement",
                "Bibliotheque de Gems par famille",
                "Cible : 10 pratiques en 3 mois",
            ],
        ),
        (
            PYL_TEAL_BLUE,
            "Accompagner les Medium",
            [
                f"{tiers.get('Medium', 0)} UC — sprints de 4 semaines",
                "Champion referent par cluster",
                "NotebookLM et App Script comme vecteurs",
            ],
        ),
        (
            PYL_DANGER,
            f"Escalader {it_count} cas IT",
            [
                "Registre de dependances systèmes",
                "Atelier de priorisation IT",
                "Zero prototype SFDC/SAP sans validation",
            ],
        ),
        (
            PYL_NAVY,
            "Creer un track Data Champion",
            [
                "F7 — famille la plus dense et complexe",
                "Profil Python/BigQuery requis",
                "Certification Data Champion a creer",
            ],
        ),
    ]

    CARD_W = Inches(3.0)
    CARD_H = Inches(3.8)

    for i, (color, title, bullets) in enumerate(recs):
        col = i % 2
        row_idx = i // 2
        left = Inches(0.5) + col * (CARD_W + Inches(0.25))
        top = Inches(1.85) + row_idx * (CARD_H + Inches(0.2))
        add_content_card(slide, left, top, CARD_W, CARD_H, color, title, bullets)

    # Roadmap droite
    roadmap_items = [
        ("M+1", "Atelier priorisation Quick Wins par cluster"),
        ("M+3", "Bibliotheque Gems V1 disponible"),
        ("M+6", "50% des Medium en production"),
        ("M+12", "Projets Large lances avec sponsors IT"),
    ]
    for j, (month, action) in enumerate(roadmap_items):
        top_r = Inches(1.85) + j * Inches(0.95)
        badge = add_rect(
            slide,
            Inches(6.9),
            top_r,
            Inches(0.8),
            Inches(0.4),
            fill=PYL_YELLOW,
        )
        bt = badge.text_frame
        bp = bt.paragraphs[0]
        bp.alignment = PP_ALIGN.CENTER
        br = bp.add_run()
        br.text = month
        br.font.name = "Poppins"
        br.font.bold = True
        br.font.size = Pt(10)
        br.font.color.rgb = PYL_NAVY_DARK

        add_textbox(
            slide,
            Inches(7.8),
            top_r,
            Inches(4.9),
            Inches(0.5),
            action,
            font_size=11,
            color=PYL_BODY_GREY,
        )


def build_best_practices_slide(prs: Presentation, page: int) -> None:
    """Slide Bonnes pratiques AI Champions."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(
        slide,
        "11",
        "Bonnes pratiques AI Champions",
        "Recommandations methodologiques pour perenniser les realisations",
        page,
    )

    practices = [
        (
            PYL_NAVY,
            "Gestion des Prompts",
            [
                "Stocker dans un Google Doc dédié",
                "Versionner : Prompt_V1.0_AAAA-MM-JJ",
                "Documenter contexte et exemples",
                "Centraliser : Drive AI Champions",
            ],
        ),
        (
            PYL_TEAL_BLUE,
            "Discipline App Script",
            [
                "1 script = 1 fichier nomme",
                "Commenter les blocs principaux",
                "Pas de secrets en clair dans le code",
                "README + tests sur données fictives",
            ],
        ),
        (
            PYL_TEAL,
            "Kit Champion",
            [
                "README.md : description + how-to",
                "Prompts_Vx.x.md : historique",
                "Script_Vx.x.gs : copie du code",
                "Tests.md + Changelog.md",
            ],
        ),
        (
            PYL_DANGER,
            "Quand escalader vers l'IT ?",
            [
                "Connexion systeme enterprise",
                "Donnees sensibles (RH, finance)",
                "Script instable / 10+ utilisateurs",
                "Besoin de fiabilite 24/7",
            ],
        ),
    ]

    CARD_W = Inches(3.0)
    CARD_H = Inches(4.2)

    for i, (color, title, bullets) in enumerate(practices):
        left = Inches(0.5) + i * (CARD_W + Inches(0.15))
        add_content_card(
            slide, left, Inches(1.8), CARD_W, CARD_H, color, title, bullets
        )


def build_closing_slide(prs: Presentation) -> None:
    """Slide de cloture."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=PYL_NAVY_DARK)

    add_textbox(
        slide,
        Inches(0.5),
        Inches(0.3),
        Inches(2),
        Inches(0.45),
        "Pyl.Tech",
        font_size=16,
        bold=True,
        color=PYL_YELLOW,
    )

    add_textbox(
        slide,
        Inches(1.5),
        Inches(2.5),
        Inches(10.3),
        Inches(1.2),
        "Transformez avec Pyl.Tech.",
        font_size=40,
        bold=True,
        color=PYL_WHITE,
        align=PP_ALIGN.CENTER,
    )

    add_textbox(
        slide,
        Inches(1.5),
        Inches(4.0),
        Inches(10.3),
        Inches(0.5),
        "Google Cloud Partner  |  IA  |  Data  |  Workspace",
        font_size=16,
        color=PYL_YELLOW,
        align=PP_ALIGN.CENTER,
    )

    add_textbox(
        slide,
        Inches(0.3),
        Inches(7.1),
        Inches(7),
        Inches(0.25),
        f"(c) Copyright {YEAR} Pyl.Tech",
        font_size=9,
        color=PYL_TEAL_BLUE,
    )


# ── Point d'entree ────────────────────────────────────────────────────────────


def main() -> None:
    print(f"Lecture du catalogue : {CATALOG_FILE}")  # noqa
    if not CATALOG_FILE.exists():
        print(f"Catalogue introuvable : {CATALOG_FILE}")  # noqa
        print("   Lancez d'abord : python3 src/generate_catalog.py")  # noqa
        sys.exit(1)

    df = pd.read_excel(CATALOG_FILE, sheet_name="Catalogue", header=0)
    print(f"   {len(df)} use cases charges")  # noqa

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    page = 1
    build_cover(prs)
    page += 1
    build_sommaire(prs, page)
    page += 1
    build_context_slide(prs, df, page)
    page += 1
    build_methodology_slide(prs, page)
    page += 1
    build_framework_slide(prs, page)
    page += 1
    build_kpi_slide(prs, df, page)
    page += 1
    build_complexity_slide(prs, df, page)
    page += 1
    build_families_slide(prs, df, page)
    page += 1
    build_clusters_slide(prs, df, page)
    page += 1
    build_it_slide(prs, df, page)
    page += 1
    build_architecture_slide(prs, page)
    page += 1
    build_quickwins_slide(prs, df, page)
    page += 1
    build_recommendations_slide(prs, df, page)
    page += 1
    build_best_practices_slide(prs, page)
    page += 1
    build_closing_slide(prs)

    OUTPUT_FILE.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_FILE)
    print(f"\nPresentation PowerPoint generee : {OUTPUT_FILE}")  # noqa
    print(f"   {len(prs.slides)} slides  |  {OUTPUT_FILE.stat().st_size // 1024} KB")  # noqa
    print(f"\n   Slides :")  # noqa
    slides_info = [
        "01 — Couverture",
        "02 — Sommaire",
        "03 — Contexte de la mission",
        "04 — Méthodologie (6 phases)",
        "05 — Cadre d'analyse (framework scoring)",
        "06 — Chiffres clés (6 KPI cards)",
        "07 — Répartition par complexité (3 cards)",
        "08 — Familles fonctionnelles",
        "09 — Répartition par cluster",
        "10 — Points d'attention IT",
        "11 — Architectures de reference (3 tiers)",
        "12 — Top Quick Wins",
        "13 — Recommandations & Roadmap",
        "14 — Bonnes pratiques AI Champions",
        "15 — Closing",
    ]
    for info in slides_info:
        print(f"   . {info}")  # noqa


if __name__ == "__main__":
    main()
